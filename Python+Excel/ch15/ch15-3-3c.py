from pptx import Presentation

prs = Presentation("自動化建立PPT簡報5.pptx")
slide = prs.slides[1]
tf = slide.shapes[1].text_frame
paragraph = tf.paragraphs[1]
print(paragraph.text)
print(tf.paragraphs[2].text)