from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

prs = Presentation()
# 新增只有標題的投影片
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "自動化在PPT投影片繪製圖表"

chart_data = CategoryChartData()
chart_data.categories = ['國文', '英文', '數學']
chart_data.add_series('陳會安', (89, 76, 82))
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
)

prs.save("自動化在PPT投影片繪製圖表.pptx")