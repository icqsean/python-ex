from pptx import Presentation

prs = Presentation("自動化建立PPT簡報5.pptx")
slide = prs.slides[1]
tf = slide.shapes[1].text_frame
for idx, paragraph in enumerate(tf.paragraphs):
    print("文字段落:", idx, "-", paragraph.text)
   