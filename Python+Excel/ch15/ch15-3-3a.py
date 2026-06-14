from pptx import Presentation

prs = Presentation("自動化建立PPT簡報5.pptx")
slide = prs.slides[1]
for idx, shape in enumerate(slide.shapes):
    if shape.has_text_frame:
        print("文字框:", idx, "-", shape.text_frame)
        print(shape.text_frame.text)
