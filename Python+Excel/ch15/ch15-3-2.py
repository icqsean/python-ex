from pptx import Presentation

prs = Presentation()
# 新增標題頁投影片
title_slide_layout = prs.slide_layouts[0] 
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
# title = slide.placeholders[0]
title.text = "自動化建立PPT簡報"
subtitle = slide.placeholders[1]
subtitle.text = "作者：陳會安"
prs.save("自動化建立PPT簡報.pptx")