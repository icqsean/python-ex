from pptx import Presentation

prs = Presentation("自動化建立PPT簡報5.pptx")
print("投影片數:", len(prs.slides))
for idx, slide in enumerate(prs.slides):
    print("投影片:", idx, "-", slide)
