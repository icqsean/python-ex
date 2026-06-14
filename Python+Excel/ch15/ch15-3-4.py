from pptx import Presentation

prs = Presentation("自動化建立PPT簡報5.pptx")
slides = list(prs.slides._sldIdLst)
idx = 3   # 刪除的投影片索引, 從0開始
prs.slides._sldIdLst.remove(slides[idx])
prs.save("自動化PPT簡報的投影片管理_刪除.pptx")