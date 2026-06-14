from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation("自動化建立PPT簡報4.pptx")
# 新增完全空白的投影片
blank_slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank_slide_layout)
left = top = width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "自動化建立PPT表格"
tf.paragraphs[0].font.size = Pt(30)
rows = cols = 2
left = top = Inches(2.0)
width = Inches(6.0)
height = Inches(0.8)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table
# 指定欄位寬度
table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(4.0)
# 寫入欄位標題
table.cell(0, 0).text = "Word文件"
table.cell(0, 1).text = "PowerPoint簡報"
# 寫入資料列
table.cell(1, 0).text = 'python-docx'
table.cell(1, 1).text = 'python-pptx'
prs.save("自動化建立PPT簡報5.pptx")