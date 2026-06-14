from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches

prs = Presentation()
# 新增只有標題的投影片
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "自動化在PPT投影片繪製圖表"

chart_data = CategoryChartData()
chart_data.categories = ['陳會安', '江小魚', '王陽明']
chart_data.add_series('國文', (89, 78, 75))
chart_data.add_series('英文', (76, 90, 66))
chart_data.add_series('數學', (82, 76, 66))
x, y, cx, cy = Inches(2), Inches(2), Inches(6), Inches(4.5)
chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.RIGHT
chart.legend.include_in_layout = False

prs.save("自動化在PPT投影片繪製圖表2.pptx")