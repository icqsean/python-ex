from pptx import Presentation

prs = Presentation("自動化建立PPT簡報5.pptx")
text_paras = []

for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text_paras.append(paragraph.text)
                
print(text_paras)