from pptx import Presentation
from pptx.util import Inches, Pt

img_path = "penguins.png"
prs = Presentation("自動化建立PPT簡報3.pptx")
# 新增只有標題的投影片
title_only_slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(title_only_slide_layout)
title = slide.shapes.title
title.text = "自動化建立PPT圖片"
left = top = Inches(2)
height = Inches(4)
pic = slide.shapes.add_picture(img_path, left, top, height=height)
left = Inches(5)
height = Inches(4)
pic = slide.shapes.add_picture(img_path, left, top, height=height)
prs.save("自動化建立PPT簡報4.pptx")