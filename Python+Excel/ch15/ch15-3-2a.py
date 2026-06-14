from pptx import Presentation

prs = Presentation("自動化建立PPT簡報.pptx")
# 新增標題和內容頁投影片
title_content_slide_layout = prs.slide_layouts[1] 
slide = prs.slides.add_slide(title_content_slide_layout)
title = slide.shapes.title
title.text = "自動化建立文字內容和清單"
body_shape = slide.placeholders[1]
tf = body_shape.text_frame
tf.text = "Python辦公室自動化套件"
p = tf.add_paragraph()
p.text = "Word自動化"
p.level = 1
p = tf.add_paragraph()
p.text = "python-docx套件"
p.level = 2
p = tf.add_paragraph()
p.text = "Excel自動化"
p.level = 1
p = tf.add_paragraph()
p.text = "openpyxl套件"
p.level = 2
prs.save("自動化建立PPT簡報2.pptx")