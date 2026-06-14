from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation("自動化建立PPT簡報2.pptx")
# 編輯第2頁投影片
slide = prs.slides[1]
left = top = width = height = Inches(12/2.54)
txBox = slide.shapes.add_textbox(left, top, width, height)
tf = txBox.text_frame
tf.text = "PowerPoint自動化"
tf.paragraphs[0].font.size = Pt(30)
p = tf.add_paragraph()
p.text = "python-pptx套件"
p.font.bold = True
p = tf.add_paragraph()
p.text = "python-pptx套件"
p.font.italic = True
p.font.color.rgb = RGBColor(247, 150, 70)
prs.save("自動化建立PPT簡報3.pptx")