from pptx import Presentation

prs = Presentation("自動化建立PPT簡報5.pptx")
slides = list(prs.slides._sldIdLst)
old_idx = 3   # 原投影片索引, 從0開始
new_idx = 2   # 新投影片索引, 從0開始
prs.slides._sldIdLst.remove(slides[old_idx])
prs.slides._sldIdLst.insert(new_idx, slides[old_idx])
prs.save("自動化PPT簡報的投影片管理_調整順序.pptx")